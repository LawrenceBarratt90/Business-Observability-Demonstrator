#!/usr/bin/env python3
"""
Business Observability Forge — Partner Event PowerPoint Generator
Generates a polished, visual presentation for the partner event.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
import os
import math

# ═══════════════════════════════════════════════════════════════
# BRAND COLORS
# ═══════════════════════════════════════════════════════════════
DT_GREEN     = RGBColor(0x73, 0xBE, 0x28)  # Dynatrace green
DT_DARK      = RGBColor(0x14, 0x16, 0x31)  # Deep navy
DT_PURPLE    = RGBColor(0x6F, 0x2D, 0xA8)  # Accent purple
DT_BLUE      = RGBColor(0x14, 0x96, 0xFF)  # Bright blue
DT_CYAN      = RGBColor(0x00, 0xB4, 0xD8)  # Cyan accent
DT_ORANGE    = RGBColor(0xFF, 0x6B, 0x35)  # Warm orange
DT_RED       = RGBColor(0xE8, 0x3E, 0x3E)  # Alert red
DT_YELLOW    = RGBColor(0xFF, 0xC1, 0x07)  # Warning yellow
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT   = RGBColor(0xE0, 0xE0, 0xE0)
GRAY_MID     = RGBColor(0x9E, 0x9E, 0x9E)
GRAY_DARK    = RGBColor(0x42, 0x42, 0x42)
BG_DARK      = RGBColor(0x0D, 0x0F, 0x1C)  # Slide background
BG_CARD      = RGBColor(0x1A, 0x1D, 0x33)  # Card background
BG_CARD_LT   = RGBColor(0x22, 0x25, 0x3D)  # Lighter card

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ═══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════

def add_dark_bg(slide, color=BG_DARK):
    """Add a full-slide dark background rectangle."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # Send to back
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return bg

def add_accent_bar(slide, x, y, w, h, color=DT_GREEN):
    """Add a small colored accent bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def add_gradient_bar(slide, y=0, height=Inches(0.06), color=DT_GREEN):
    """Add a thin accent line across the top."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Segoe UI', line_spacing=1.2):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_rich_text_box(slide, left, top, width, height, paragraphs_data,
                      font_name='Segoe UI'):
    """Add text box with multiple styled paragraphs.
    paragraphs_data: list of dicts with keys: text, size, color, bold, align, spacing_after
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pd in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pd.get('text', '')
        p.font.size = Pt(pd.get('size', 18))
        p.font.color.rgb = pd.get('color', WHITE)
        p.font.bold = pd.get('bold', False)
        p.font.name = font_name
        p.alignment = pd.get('align', PP_ALIGN.LEFT)
        p.space_after = Pt(pd.get('spacing_after', 6))
        if pd.get('italic'):
            p.font.italic = True
    return txBox

def add_card(slide, left, top, width, height, color=BG_CARD, corner=Inches(0.15)):
    """Add a rounded-corner card shape."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    return card

def add_icon_circle(slide, left, top, size, color, text='', font_size=20):
    """Add a circle with icon text (emoji or symbol)."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    if text:
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return circle

def add_stat_card(slide, left, top, width, height, number, label, accent_color=DT_GREEN):
    """Add a stat card with a big number and label."""
    card = add_card(slide, left, top, width, height, BG_CARD)
    add_accent_bar(slide, left, top, width, Inches(0.05), accent_color)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.25), width - Inches(0.5), Inches(0.8),
                 number, font_size=36, color=accent_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.95), width - Inches(0.3), Inches(0.6),
                 label, font_size=13, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
    return card

def add_section_header(slide, number, title, color=DT_GREEN):
    """Add section number + title at top of content slides."""
    add_gradient_bar(slide, 0, Inches(0.04), color)
    # Section number circle
    circle = add_icon_circle(slide, Inches(0.6), Inches(0.35), Inches(0.55), color, str(number), 22)
    # Title
    add_text_box(slide, Inches(1.35), Inches(0.3), Inches(10), Inches(0.65),
                 title, font_size=32, color=WHITE, bold=True)
    return circle

def make_chevron_flow(slide, items, y_start, colors=None):
    """Create a horizontal chevron/arrow flow of items."""
    n = len(items)
    total_w = Inches(11.5)
    item_w = total_w / n
    start_x = Inches(0.9)
    if not colors:
        colors = [DT_GREEN, DT_BLUE, DT_CYAN, DT_PURPLE, DT_ORANGE, DT_GREEN]
    for i, item in enumerate(items):
        c = colors[i % len(colors)]
        x = start_x + item_w * i
        # Chevron shape
        chev = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, int(x), int(y_start),
                                       int(item_w - Inches(0.08)), Inches(0.9))
        chev.fill.solid()
        chev.fill.fore_color.rgb = c
        chev.line.fill.background()
        tf = chev.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Segoe UI'

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_bg(slide, RGBColor(0x08, 0x0A, 0x18))

# Decorative geometric shapes
for i in range(6):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
        Inches(9 + i * 0.7), Inches(0.5 + i * 0.9),
        Inches(1.5 - i * 0.15), Inches(1.5 - i * 0.15))
    s.fill.solid()
    s.fill.fore_color.rgb = DT_GREEN
    s.fill.fore_color.brightness = -0.6 + i * 0.08
    s.line.fill.background()
    s.rotation = i * 15

# Top accent
add_gradient_bar(slide, 0, Inches(0.05), DT_GREEN)

# Main title
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(9), Inches(0.5),
             'BUSINESS OBSERVABILITY', font_size=22, color=DT_GREEN, bold=True,
             font_name='Segoe UI Semibold')

add_text_box(slide, Inches(0.8), Inches(1.85), Inches(11), Inches(1.4),
             'FORGE', font_size=72, color=WHITE, bold=True,
             font_name='Segoe UI Black')

add_accent_bar(slide, Inches(0.8), Inches(3.3), Inches(2), Inches(0.06), DT_GREEN)

add_text_box(slide, Inches(0.8), Inches(3.6), Inches(10), Inches(0.8),
             'AI-Powered Customer Journey Simulation\nfor Dynatrace Partners',
             font_size=24, color=GRAY_LIGHT, font_name='Segoe UI Light')

# Stats bar at bottom
bar_y = Inches(5.3)
add_card(slide, Inches(0.5), bar_y, Inches(12.3), Inches(1.5), BG_CARD)

stats = [
    ('55', 'Industry\nVerticals', DT_GREEN),
    ('256+', 'Journey\nTemplates', DT_BLUE),
    ('4', 'AI\nAgents', DT_PURPLE),
    ('< 15 min', 'Setup\nTime', DT_CYAN),
    ('v2.22.5', 'Current\nVersion', DT_ORANGE),
]
for i, (num, lbl, col) in enumerate(stats):
    x = Inches(0.8 + i * 2.4)
    add_text_box(slide, x, bar_y + Inches(0.15), Inches(2), Inches(0.7),
                 num, font_size=32, color=col, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, bar_y + Inches(0.8), Inches(2), Inches(0.6),
                 lbl, font_size=12, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# Date
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(6), Inches(0.4),
             'Partner Event  |  April 2026  |  Dynatrace', font_size=14,
             color=GRAY_MID, font_name='Segoe UI Light')


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 1, 'The Partner Challenge')

# Three problem cards
problems = [
    ('Every Customer\nis Different', 
     'A bank cares about fraud rates.\nA hospital cares about triage accuracy.\nA retailer cares about cart abandonment.\n\nYou can\'t build a demo for every vertical.',
     DT_ORANGE, '!'),
    ('Demo Setup\nTakes Weeks',
     'Real services, real data, real dashboards.\nMost partners resort to PowerPoint\nand say "imagine this..."\n\nThat doesn\'t close deals.',
     DT_RED, '~'),
    ('C-Suite Doesn\'t\nSpeak MELT',
     'They speak revenue, risk, cost,\nand customer experience.\n\nWe need to translate\nobservability into their language.',
     DT_PURPLE, '$'),
]

for i, (title, desc, color, icon) in enumerate(problems):
    x = Inches(0.7 + i * 4.1)
    y = Inches(1.5)
    card = add_card(slide, x, y, Inches(3.8), Inches(5.2), BG_CARD)
    add_accent_bar(slide, x, y, Inches(3.8), Inches(0.06), color)
    add_icon_circle(slide, x + Inches(1.45), y + Inches(0.35), Inches(0.7), color, icon, 28)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.3), Inches(3.2), Inches(0.9),
                 title, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(2.4), Inches(3.2), Inches(2.6),
                 desc, font_size=14, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: WHAT IS THE FORGE — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 2, 'What Is the Business Observability Forge?')

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.6),
             'An AI-powered customer journey simulation engine that generates real, instrumented '
             'microservices within Dynatrace — in minutes, not weeks.',
             font_size=16, color=GRAY_LIGHT, font_name='Segoe UI Light')

# Two-part architecture
# LEFT: Engine
card_w = Inches(5.5)
card_h = Inches(4.5)

# Engine card
add_card(slide, Inches(0.7), Inches(2.1), card_w, card_h, BG_CARD)
add_accent_bar(slide, Inches(0.7), Inches(2.1), card_w, Inches(0.06), DT_GREEN)
add_icon_circle(slide, Inches(1.0), Inches(2.4), Inches(0.5), DT_GREEN, 'E', 18)
add_text_box(slide, Inches(1.65), Inches(2.35), Inches(4), Inches(0.45),
             'THE ENGINE', font_size=20, color=DT_GREEN, bold=True)
add_text_box(slide, Inches(1.65), Inches(2.75), Inches(4), Inches(0.35),
             'Node.js Server  |  Your VM / EC2 / Laptop', font_size=11, color=GRAY_MID)

engine_features = [
    'Dynamic child service spawning (ports 8081-8200)',
    'Real HTTP microservices with trace propagation',
    'OneAgent auto-detection & Smartscape topology',
    'Business event emission at every journey step',
    'AI agent framework (Nemesis, Fix-It, Librarian)',
    'Feature flag management per service',
    'Auto-load traffic generation (30-60 journeys/min)',
    '256+ pre-built journey templates',
]
for j, feat in enumerate(engine_features):
    add_text_box(slide, Inches(1.3), Inches(3.25 + j * 0.38), Inches(4.5), Inches(0.35),
                 f'  {feat}', font_size=12, color=GRAY_LIGHT)
    # Green dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.05), Inches(3.32 + j * 0.38),
                                  Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = DT_GREEN
    dot.line.fill.background()

# Arrow between
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.4), Inches(3.8),
                                Inches(0.8), Inches(0.6))
arrow.fill.solid()
arrow.fill.fore_color.rgb = DT_CYAN
arrow.line.fill.background()
add_text_box(slide, Inches(6.15), Inches(4.5), Inches(1.3), Inches(0.4),
             'EdgeConnect', font_size=9, color=DT_CYAN, alignment=PP_ALIGN.CENTER)

# Forge UI card
add_card(slide, Inches(7.3), Inches(2.1), card_w, card_h, BG_CARD)
add_accent_bar(slide, Inches(7.3), Inches(2.1), card_w, Inches(0.06), DT_PURPLE)
add_icon_circle(slide, Inches(7.6), Inches(2.4), Inches(0.5), DT_PURPLE, 'F', 18)
add_text_box(slide, Inches(8.25), Inches(2.35), Inches(4), Inches(0.45),
             'THE FORGE UI', font_size=20, color=DT_PURPLE, bold=True)
add_text_box(slide, Inches(8.25), Inches(2.75), Inches(4), Inches(0.35),
             'Dynatrace AppEngine  |  Native In-Tenant', font_size=11, color=GRAY_MID)

ui_tabs = [
    ('Home', 'Journey builder & launcher'),
    ('Solutions', 'Industry vertical templates'),
    ('Dashboards', 'AI-generated Dynatrace dashboards'),
    ('Services', 'Live service monitoring'),
    ('Chaos Control', 'Nemesis agent interface'),
    ('Fix-It Agent', 'AI remediation control'),
    ('Demo Guide', 'Interactive walkthrough'),
    ('Settings', 'Credentials & configuration'),
]
for j, (tab, desc) in enumerate(ui_tabs):
    add_text_box(slide, Inches(7.9), Inches(3.25 + j * 0.38), Inches(1.5), Inches(0.35),
                 tab, font_size=12, color=DT_CYAN, bold=True)
    add_text_box(slide, Inches(9.4), Inches(3.25 + j * 0.38), Inches(3), Inches(0.35),
                 desc, font_size=11, color=GRAY_LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: HOW IT WORKS — JOURNEY FLOW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 3, 'How It Works: From Zero to Business Observability')

# Timeline steps
steps = [
    ('Pick Industry', 'Select from 55\nverticals', DT_GREEN, '1'),
    ('Choose Journey', 'Pre-built templates\nfor every domain', DT_BLUE, '2'),
    ('Services Spawn', '6-8 real microservices\nauto-created', DT_CYAN, '3'),
    ('Traces Flow', 'Distributed traces\nacross all services', DT_PURPLE, '4'),
    ('BizEvents Emit', 'Revenue & KPI data\nat every step', DT_ORANGE, '5'),
    ('Dashboards\nGenerate', 'AI builds DQL\ndashboard tiles', DT_GREEN, '6'),
]

for i, (title, desc, color, num) in enumerate(steps):
    x = Inches(0.5 + i * 2.1)
    y_base = Inches(1.6)
    
    # Number circle
    add_icon_circle(slide, x + Inches(0.65), y_base, Inches(0.6), color, num, 22)
    # Connecting line
    if i < len(steps) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            x + Inches(1.35), y_base + Inches(0.25), Inches(1.45), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = GRAY_DARK
        line.line.fill.background()
    
    # Title
    add_text_box(slide, x, y_base + Inches(0.8), Inches(2), Inches(0.6),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Desc
    add_text_box(slide, x, y_base + Inches(1.4), Inches(2), Inches(0.7),
                 desc, font_size=11, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# Demo scenario card
add_card(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.0), BG_CARD)
add_accent_bar(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.05), DT_GREEN)

add_text_box(slide, Inches(0.9), Inches(4.2), Inches(6), Inches(0.5),
             'LIVE DEMO: Retail Banking — Loan Application', font_size=18, color=DT_GREEN, bold=True)

# Journey flow chevrons
make_chevron_flow(slide, [
    'Application\nSubmission',
    'Credit\nCheck',
    'Document\nVerification', 
    'Underwriting\nDecision',
    'Loan\nApproval',
    'Funds\nDisbursement'
], Inches(4.85))

demo_points = [
    '6 independent Node.js services spawn automatically — each on its own port',
    'OneAgent detects every service — full Smartscape topology created from real HTTP traffic',
    'Each step emits business events: loan value, risk score, customer segment, conversion rate',
    'AI generates a revenue dashboard from the journey data — no DQL writing required',
]
for j, point in enumerate(demo_points):
    add_text_box(slide, Inches(1.3), Inches(5.9 + j * 0.3), Inches(10.5), Inches(0.3),
                 f'     {point}', font_size=12, color=GRAY_LIGHT)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.05), Inches(5.96 + j * 0.3),
                                  Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = DT_GREEN
    dot.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: BUSINESS EVENTS — THE REVENUE BRIDGE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 4, 'Business Events: The Revenue Bridge')

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
             'Every journey step emits structured business events — translating MELT into the language the C-suite speaks.',
             font_size=15, color=GRAY_LIGHT, font_name='Segoe UI Light')

# BizEvent JSON card
add_card(slide, Inches(0.5), Inches(1.9), Inches(5.5), Inches(5.0), RGBColor(0x12, 0x14, 0x28))
add_accent_bar(slide, Inches(0.5), Inches(1.9), Inches(5.5), Inches(0.05), DT_CYAN)
add_text_box(slide, Inches(0.8), Inches(2.05), Inches(3), Inches(0.35),
             'Sample Business Event', font_size=13, color=DT_CYAN, bold=True)

json_lines = [
    ('{', GRAY_MID),
    ('  "event.type":        "BIZ_EVENT"', GRAY_LIGHT),
    ('  "companyName":       "First National Bank"', GRAY_LIGHT),
    ('  "journeyType":       "Loan Application"', GRAY_LIGHT),
    ('  "stepName":          "CreditCheck"', GRAY_LIGHT),
    ('  "orderTotal":         45,000', DT_GREEN),
    ('  "customerLifetimeValue": 128,000', DT_GREEN),
    ('  "conversionRate":     78%', DT_ORANGE),
    ('  "riskScore":          0.23', DT_YELLOW),
    ('  "approvalTime":       4.2s', DT_CYAN),
    ('  "channel":           "mobile_app"', GRAY_LIGHT),
    ('  "customerSegment":   "Prime"', DT_PURPLE),
    ('}', GRAY_MID),
]
for j, (line, color) in enumerate(json_lines):
    add_text_box(slide, Inches(0.8), Inches(2.45 + j * 0.33), Inches(5), Inches(0.3),
                 line, font_size=11, color=color, font_name='Consolas')

# Right side: what this unlocks
add_card(slide, Inches(6.3), Inches(1.9), Inches(6.3), Inches(5.0), BG_CARD)
add_accent_bar(slide, Inches(6.3), Inches(1.9), Inches(6.3), Inches(0.05), DT_GREEN)
add_text_box(slide, Inches(6.6), Inches(2.05), Inches(5), Inches(0.35),
             'What This Unlocks', font_size=13, color=DT_GREEN, bold=True)

insights = [
    ('Revenue Impact Tracking', 
     '"Checkout latency caused £2.3K revenue loss this hour"',
     DT_GREEN),
    ('Customer Churn Correlation',
     '"23 failed logins x £240 LTV = £5.5K churn risk"',
     DT_ORANGE),
    ('AI Model Accuracy',
     '"Credit scoring 94.2% accurate — 5.8% false decline rate"',
     DT_CYAN),
    ('Conversion Funnel Analysis',
     '"78% submission-to-approval, 12% drop at document verification"',
     DT_PURPLE),
    ('SLA-to-Revenue Mapping',
     '"P99 latency breach on payment = £18K/hr in failed settlements"',
     DT_RED),
]

for j, (title, quote, color) in enumerate(insights):
    y_pos = Inches(2.5 + j * 0.9)
    add_icon_circle(slide, Inches(6.6), y_pos, Inches(0.35), color, '', 10)
    add_text_box(slide, Inches(7.1), y_pos - Inches(0.05), Inches(5), Inches(0.35),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.1), y_pos + Inches(0.3), Inches(5.2), Inches(0.35),
                 quote, font_size=11, color=GRAY_LIGHT, font_name='Segoe UI Light')


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: 55 INDUSTRY VERTICALS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 5, '55 Industry Verticals')

add_text_box(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4),
             'Pre-built journey templates, BizEvent schemas, and ROI language for every industry.',
             font_size=14, color=GRAY_LIGHT, font_name='Segoe UI Light')

# Vertical category grid
categories = [
    ('Financial Services', ['Retail Banking', 'Wealth Mgmt', 'Payments & Fintech',
                            'Insurance', 'Accounting', 'Lottery & Betting'], DT_GREEN, '6'),
    ('Healthcare', ['Healthcare & Life Sci', 'Pharmaceuticals', 'Veterinary'], DT_BLUE, '3'),
    ('Retail & Consumer', ['E-commerce', 'Fashion & Luxury', 'Food & Beverage',
                           'Beauty & Cosmetics', 'Fitness', 'Food Delivery',
                           'Hospitality', 'Gaming'], DT_PURPLE, '8'),
    ('Travel & Transport', ['Airlines', 'Automotive', 'Railway', 'Ride-hailing',
                            'Shipping', 'EV Charging', 'Logistics'], DT_CYAN, '7'),
    ('Industrial', ['Manufacturing', 'Construction', 'Mining', 'Chemical',
                    'Semiconductors', 'Defence', 'Agriculture'], DT_ORANGE, '7'),
    ('Utilities & Energy', ['Energy & Utilities', 'Water & Waste', 'Recycling'], DT_YELLOW, '3'),
    ('Technology', ['Data Centres', 'Cybersecurity', 'Telecoms', 'Smart Cities',
                    'Robotics', 'Space & Satellite'], DT_RED, '6'),
    ('Education & Prof', ['Education', 'Government', 'Consulting', 'Legal'], DT_GREEN, '4'),
    ('Media & Digital', ['Media & Entertainment', 'Publishing', 'Music & Audio',
                         'Advertising', 'Social Media', 'Marketplaces', 'Real Estate'], DT_BLUE, '7'),
    ('Human Services', ['HR & Workforce', 'Nonprofit & Health'], DT_PURPLE, '2'),
    ('Emerging', ['Environmental & ESG', 'BizObs Platform'], DT_CYAN, '2'),
]

# Layout in 4 columns, wrapping
cols = 4
col_w = Inches(3.0)
row_h_base = Inches(1.8)
gap = Inches(0.15)
start_x = Inches(0.5)
start_y = Inches(1.65)

col_items = [[], [], [], []]  # Distribute by filling columns
heights = [0, 0, 0, 0]

for cat in categories:
    # Find shortest column
    min_col = heights.index(min(heights))
    col_items[min_col].append(cat)
    item_h = 0.6 + len(cat[1]) * 0.22
    heights[min_col] += item_h

for col_idx in range(cols):
    x = start_x + col_idx * (col_w + gap)
    y_offset = start_y
    for (cat_name, items, color, count) in col_items[col_idx]:
        card_h = Inches(0.55 + len(items) * 0.22)
        add_card(slide, x, y_offset, col_w, card_h, BG_CARD)
        add_accent_bar(slide, x, y_offset, Inches(0.06), card_h, color)
        
        # Category header with count badge
        add_text_box(slide, x + Inches(0.15), y_offset + Inches(0.05), Inches(2.2), Inches(0.3),
                     cat_name, font_size=11, color=color, bold=True)
        # Count badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            x + col_w - Inches(0.55), y_offset + Inches(0.07), Inches(0.4), Inches(0.22))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = count
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        for k, item in enumerate(items):
            add_text_box(slide, x + Inches(0.2), y_offset + Inches(0.32 + k * 0.22),
                         Inches(2.6), Inches(0.22),
                         f'  {item}', font_size=9, color=GRAY_LIGHT)
        
        y_offset += card_h + gap


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: TOP VERTICALS FOR PARTNER CONVERSATIONS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 6, 'Top Verticals for Partner Conversations')

verticals = [
    ('Retail Banking', 'AI fraud detection false positive rate — money left on the table', DT_GREEN),
    ('Healthcare', 'AI triage accuracy — prove clinical decision support works', DT_BLUE),
    ('E-commerce', 'Every 100ms latency = 1% conversion loss — measure it', DT_PURPLE),
    ('Insurance', 'Straight-through processing rate = profitability driver', DT_CYAN),
    ('Telecoms', '5G slice performance directly impacts enterprise SLA revenue', DT_ORANGE),
    ('Automotive', 'OTA update success rate by vehicle model and region', DT_YELLOW),
    ('Airlines', 'Check-in abandonment at doc verification = lost ancillary revenue', DT_RED),
    ('Energy', 'Predictive maintenance catching failures before outages', DT_GREEN),
    ('Pharma', 'AI patient matching reduced Phase III enrollment 40%', DT_BLUE),
    ('Government', 'Citizen portal uptime during benefit enrollment periods', DT_PURPLE),
]

for i, (vert, hook, color) in enumerate(verticals):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.35 + row * 1.15)
    
    add_card(slide, x, y, Inches(6.0), Inches(1.0), BG_CARD)
    add_accent_bar(slide, x, y, Inches(0.06), Inches(1.0), color)
    
    # Number
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(0.4), Inches(0.35),
                 str(i + 1), font_size=18, color=color, bold=True)
    # Vertical name
    add_text_box(slide, x + Inches(0.6), y + Inches(0.1), Inches(2.5), Inches(0.35),
                 vert, font_size=16, color=WHITE, bold=True)
    # Hook
    add_text_box(slide, x + Inches(0.6), y + Inches(0.5), Inches(5), Inches(0.4),
                 f'"{hook}"', font_size=11, color=GRAY_LIGHT, font_name='Segoe UI Light')


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: CHAOS ENGINEERING & AI REMEDIATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 7, 'Chaos Engineering & AI Remediation', DT_RED)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
             'Observability isn\'t just watching things work — it\'s understanding what happens when they don\'t.',
             font_size=15, color=GRAY_LIGHT, font_name='Segoe UI Light')

# Nemesis card
add_card(slide, Inches(0.5), Inches(1.9), Inches(6.0), Inches(5.0), BG_CARD)
add_accent_bar(slide, Inches(0.5), Inches(1.9), Inches(6.0), Inches(0.06), DT_RED)
add_text_box(slide, Inches(0.8), Inches(2.1), Inches(4), Inches(0.4),
             'NEMESIS AGENT — Chaos Engineering', font_size=16, color=DT_RED, bold=True)

chaos_types = [
    ('Service Unavailable (503)', 'Blocked transactions, revenue loss', DT_RED),
    ('Timeout (504)', 'Customer abandonment, SLA breach', DT_ORANGE),
    ('Connection Refused', 'Cascading failures, queue buildup', DT_YELLOW),
    ('Internal Error (500)', 'Data inconsistency, retry storms', DT_PURPLE),
    ('Slow Response (2-10x)', 'UX degradation, conversion drop', DT_CYAN),
    ('Circuit Breaker Trip', 'Graceful degradation test', DT_BLUE),
]

for j, (chaos, impact, color) in enumerate(chaos_types):
    y_pos = Inches(2.7 + j * 0.62)
    add_card(slide, Inches(0.7), y_pos, Inches(5.6), Inches(0.52), BG_CARD_LT)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y_pos + Inches(0.15),
                                  Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text_box(slide, Inches(1.2), y_pos + Inches(0.02), Inches(2.5), Inches(0.25),
                 chaos, font_size=11, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.2), y_pos + Inches(0.26), Inches(4.8), Inches(0.22),
                 impact, font_size=10, color=GRAY_MID)

# Fix-It card
add_card(slide, Inches(6.8), Inches(1.9), Inches(5.8), Inches(5.0), BG_CARD)
add_accent_bar(slide, Inches(6.8), Inches(1.9), Inches(5.8), Inches(0.06), DT_GREEN)
add_text_box(slide, Inches(7.1), Inches(2.1), Inches(4), Inches(0.4),
             'FIX-IT AGENT — AI Remediation', font_size=16, color=DT_GREEN, bold=True)

# Flow diagram
fix_steps = [
    ('1  DETECT', 'Davis AI problem detection fires', DT_RED),
    ('2  CORRELATE', 'Maps problem to business impact\n(which journey, which step, how much £)', DT_ORANGE),
    ('3  REMEDIATE', 'Auto-triggers workflow: restart service,\nroll back feature flag, or escalate', DT_GREEN),
    ('4  AUDIT', 'Librarian logs everything:\nwhat happened, why, business cost', DT_BLUE),
]

for j, (step, desc, color) in enumerate(fix_steps):
    y_pos = Inches(2.7 + j * 1.0)
    add_card(slide, Inches(7.0), y_pos, Inches(5.4), Inches(0.85), BG_CARD_LT)
    add_accent_bar(slide, Inches(7.0), y_pos, Inches(0.06), Inches(0.85), color)
    add_text_box(slide, Inches(7.2), y_pos + Inches(0.05), Inches(2), Inches(0.3),
                 step, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(7.2), y_pos + Inches(0.35), Inches(5), Inches(0.45),
                 desc, font_size=11, color=GRAY_LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: SOLUTIONS DEEP DIVE — Healthcare, Banking, Pharma
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 8, 'Solutions: Pre-Built Industry Integrations')

solutions = [
    {
        'title': 'Healthcare & Life Sciences',
        'tagline': '"Prove AI is improving patient outcomes — not just generating cost"',
        'integrations': ['Epic EHR', 'Cerner (Oracle Health)', 'Meditech'],
        'kpis': [
            'AI recommendation acceptance vs override rate',
            'MyChart AI triage confidence & accuracy',
            'Predictive sepsis early warning effectiveness',
            'FHIR API response times for AI consumers',
        ],
        'roi': '£180K prevented in adverse outcomes/month',
        'color': DT_BLUE,
        'icon': '+',
    },
    {
        'title': 'Retail Banking',
        'tagline': '"Prove AI is driving revenue and reducing risk — not just hype"',
        'integrations': ['FIS', 'Temenos', 'Finastra'],
        'kpis': [
            'AI fraud false positive rate & blocked revenue',
            'Credit decision accuracy vs 90-day defaults',
            'Payment routing approval rate lift %',
            'Chatbot containment rate & escalation patterns',
        ],
        'roi': '£8.4M fraud prevented / £2.3M recovered',
        'color': DT_GREEN,
        'icon': '$',
    },
    {
        'title': 'Pharmaceuticals',
        'tagline': '"Prove AI is accelerating discovery and reducing compliance risk"',
        'integrations': ['IQVIA', 'SAP S/4HANA', 'Regulatory Systems'],
        'kpis': [
            'AI patient matching precision/recall',
            'Demand forecast accuracy (MAPE) by product',
            'Batch release automation time savings',
            'Cold chain optimisation cost savings',
        ],
        'roi': '£28M+ extended patent exclusivity value',
        'color': DT_PURPLE,
        'icon': 'R',
    },
]

for i, sol in enumerate(solutions):
    x = Inches(0.4 + i * 4.2)
    y = Inches(1.35)
    card_w = Inches(4.0)
    card_h = Inches(5.8)
    
    add_card(slide, x, y, card_w, card_h, BG_CARD)
    add_accent_bar(slide, x, y, card_w, Inches(0.06), sol['color'])
    
    # Icon + Title
    add_icon_circle(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.45), sol['color'], sol['icon'], 18)
    add_text_box(slide, x + Inches(0.75), y + Inches(0.25), Inches(3), Inches(0.4),
                 sol['title'], font_size=15, color=WHITE, bold=True)
    
    # Tagline
    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.6), Inches(0.5),
                 sol['tagline'], font_size=10, color=sol['color'], font_name='Segoe UI Light')
    
    # Integrations
    add_text_box(slide, x + Inches(0.2), y + Inches(1.35), Inches(3), Inches(0.25),
                 'PRE-BUILT INTEGRATIONS', font_size=9, color=GRAY_MID, bold=True)
    for j, integ in enumerate(sol['integrations']):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.2 + j * 1.15), y + Inches(1.6), Inches(1.1), Inches(0.25))
        badge.fill.solid()
        badge.fill.fore_color.rgb = BG_CARD_LT
        badge.line.color.rgb = sol['color']
        badge.line.width = Pt(1)
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = integ
        p.font.size = Pt(8)
        p.font.color.rgb = GRAY_LIGHT
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Segoe UI'
    
    # KPIs
    add_text_box(slide, x + Inches(0.2), y + Inches(2.1), Inches(3), Inches(0.25),
                 'KEY DEMO KPIs', font_size=9, color=GRAY_MID, bold=True)
    for j, kpi in enumerate(sol['kpis']):
        add_text_box(slide, x + Inches(0.2), y + Inches(2.4 + j * 0.32), Inches(3.5), Inches(0.3),
                     f'  {kpi}', font_size=10, color=GRAY_LIGHT)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(2.47 + j * 0.32),
                                      Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = sol['color']
        dot.line.fill.background()
    
    # ROI banner
    roi_y = y + Inches(3.85)
    roi_card = add_card(slide, x + Inches(0.15), roi_y, Inches(3.7), Inches(0.7), sol['color'])
    add_text_box(slide, x + Inches(0.3), roi_y + Inches(0.03), Inches(3.4), Inches(0.2),
                 'ROI IMPACT', font_size=8, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), roi_y + Inches(0.25), Inches(3.5), Inches(0.35),
                 sol['roi'], font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: MANAGED COMPATIBILITY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 9, 'What Works on Dynatrace Managed')

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
             'The Forge works on Managed TODAY for core observability — and becomes the migration conversation for SaaS.',
             font_size=15, color=GRAY_LIGHT, font_name='Segoe UI Light')

# Works on Managed
add_card(slide, Inches(0.5), Inches(1.9), Inches(5.8), Inches(5.0), BG_CARD)
add_accent_bar(slide, Inches(0.5), Inches(1.9), Inches(5.8), Inches(0.06), DT_GREEN)
add_text_box(slide, Inches(0.8), Inches(2.05), Inches(5), Inches(0.4),
             'WORKS ON MANAGED TODAY', font_size=18, color=DT_GREEN, bold=True)

managed_features = [
    'Real microservice generation & detection',
    'Full Smartscape topology from HTTP traffic',
    'Distributed tracing across all services',
    'Davis AI problem detection',
    'Chaos engineering (feature flag injection)',
    'OneAgent service splitting (DT_TAGS)',
    'Custom events via Events API v2',
    'Executive PDF export',
    '256+ journey templates',
    'Multiple concurrent companies',
]
for j, feat in enumerate(managed_features):
    y_pos = Inches(2.6 + j * 0.42)
    # Green check
    check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos + Inches(0.05),
                                    Inches(0.2), Inches(0.2))
    check.fill.solid()
    check.fill.fore_color.rgb = DT_GREEN
    check.line.fill.background()
    tf = check.text_frame
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    
    add_text_box(slide, Inches(1.15), y_pos, Inches(4.8), Inches(0.35),
                 feat, font_size=13, color=WHITE)

# Enhanced on SaaS
add_card(slide, Inches(6.8), Inches(1.9), Inches(5.8), Inches(5.0), BG_CARD)
add_accent_bar(slide, Inches(6.8), Inches(1.9), Inches(5.8), Inches(0.06), DT_PURPLE)
add_text_box(slide, Inches(7.1), Inches(2.05), Inches(5), Inches(0.4),
             'ENHANCED ON SAAS / GRAIL', font_size=18, color=DT_PURPLE, bold=True)

saas_features = [
    ('Business event querying', 'fetch bizevents via DQL'),
    ('AI dashboard generation', 'Auto-built DQL tile dashboards'),
    ('Forge UI (AppEngine)', 'Native in-tenant control plane'),
    ('DQL notebooks', 'Ad-hoc business analysis'),
    ('Davis AI + Grail', 'Enhanced root cause analysis'),
]
for j, (feat, desc) in enumerate(saas_features):
    y_pos = Inches(2.6 + j * 0.75)
    add_card(slide, Inches(7.0), y_pos, Inches(5.4), Inches(0.65), BG_CARD_LT)
    add_accent_bar(slide, Inches(7.0), y_pos, Inches(0.06), Inches(0.65), DT_PURPLE)
    # Star icon
    star = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), y_pos + Inches(0.1),
                                   Inches(0.2), Inches(0.2))
    star.fill.solid()
    star.fill.fore_color.rgb = DT_PURPLE
    star.line.fill.background()
    add_text_box(slide, Inches(7.55), y_pos + Inches(0.03), Inches(4.5), Inches(0.3),
                 feat, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.55), y_pos + Inches(0.33), Inches(4.5), Inches(0.25),
                 desc, font_size=11, color=GRAY_MID)

# Migration message
add_card(slide, Inches(6.8), Inches(6.4), Inches(5.8), Inches(0.5), DT_PURPLE)
add_text_box(slide, Inches(7.0), Inches(6.4), Inches(5.4), Inches(0.5),
             'The Forge is a SaaS migration accelerator — not a limitation.',
             font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: AI AGENTS OVERVIEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 10, 'Four AI Agents Working Together')

agents = [
    {
        'name': 'NEMESIS',
        'subtitle': 'Chaos Engineering',
        'desc': 'Autonomous, AI-driven chaos injection via feature flags. Per-service targeting with safety locks. LLM-assisted recipe selection.',
        'color': DT_RED,
        'icon': 'N',
    },
    {
        'name': 'FIX-IT',
        'subtitle': 'AI Remediation',
        'desc': 'Detects problems via Davis AI. Correlates with business impact (revenue, CX). Triggers automated remediation workflows.',
        'color': DT_GREEN,
        'icon': 'F',
    },
    {
        'name': 'LIBRARIAN',
        'subtitle': 'Memory & Audit',
        'desc': 'Central memory store for all chaos events and AI decisions. Enables "why was this triggered?" playback. Full audit trail.',
        'color': DT_BLUE,
        'icon': 'L',
    },
    {
        'name': 'DASHBOARD',
        'subtitle': 'BI Generation',
        'desc': 'Generates and deploys Dynatrace dashboards automatically. DQL-powered tiles correlated to business journeys. Ollama LLM integration.',
        'color': DT_PURPLE,
        'icon': 'D',
    },
]

for i, agent in enumerate(agents):
    x = Inches(0.4 + i * 3.2)
    y = Inches(1.4)
    card_w = Inches(3.0)
    card_h = Inches(5.5)
    
    add_card(slide, x, y, card_w, card_h, BG_CARD)
    add_accent_bar(slide, x, y, card_w, Inches(0.06), agent['color'])
    
    # Big icon circle
    add_icon_circle(slide, x + Inches(0.95), y + Inches(0.4), Inches(0.9), agent['color'], agent['icon'], 32)
    
    # Name
    add_text_box(slide, x + Inches(0.15), y + Inches(1.55), card_w - Inches(0.3), Inches(0.4),
                 agent['name'], font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Subtitle
    add_text_box(slide, x + Inches(0.15), y + Inches(1.95), card_w - Inches(0.3), Inches(0.35),
                 agent['subtitle'], font_size=12, color=agent['color'], alignment=PP_ALIGN.CENTER)
    
    # Divider line
    add_accent_bar(slide, x + Inches(0.8), y + Inches(2.4), Inches(1.4), Inches(0.02), agent['color'])
    
    # Description
    add_text_box(slide, x + Inches(0.2), y + Inches(2.65), card_w - Inches(0.4), Inches(2.5),
                 agent['desc'], font_size=12, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# Connection arrows between agents
for i in range(3):
    x = Inches(3.25 + i * 3.2)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(3.8), Inches(0.5), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY_DARK
    arrow.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: EXECUTIVE PDF & EXPORT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 11, 'Leave-Behind: Executive PDF Export')

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
             'The Forge generates polished, C-suite-ready executive summaries — your meeting leave-behind that closes deals.',
             font_size=15, color=GRAY_LIGHT, font_name='Segoe UI Light')

# PDF preview mock
add_card(slide, Inches(0.5), Inches(1.9), Inches(5.0), Inches(5.0), WHITE)

# Mock PDF content
add_text_box(slide, Inches(0.8), Inches(2.1), Inches(4.5), Inches(0.3),
             'EXECUTIVE SUMMARY', font_size=14, color=DT_DARK, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(2.45), Inches(1.5), Inches(0.04), DT_GREEN)
add_text_box(slide, Inches(0.8), Inches(2.6), Inches(4.3), Inches(0.25),
             'First National Bank — Loan Application Journey', font_size=11, color=GRAY_DARK, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.9), Inches(4.3), Inches(1.2),
             'Key Findings:\n'
             '• AI credit decisioning reduced false decline rates by 12%\n'
             '• Recovered £2.3M in previously blocked applications\n'
             '• Fraud detection maintained 99.7% accuracy\n'
             '• P99 latency at 4.2s (target: 5s) — within SLA',
             font_size=9, color=GRAY_DARK)
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(4.3), Inches(0.8),
             'ROI Projection:\n'
             '• £8.4M fraud prevention quarterly\n'
             '• 23% faster loan processing time\n'
             '• NPS improvement: 42 → 58 projected',
             font_size=9, color=GRAY_DARK)
# Mock chart area
chart_area = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(5.1), Inches(4.3), Inches(1.5))
chart_area.fill.solid()
chart_area.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
chart_area.line.color.rgb = GRAY_LIGHT
chart_area.line.width = Pt(0.5)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(3), Inches(0.4),
             '[Journey Flow Visualisation]', font_size=10, color=GRAY_MID, alignment=PP_ALIGN.CENTER)

# Right side: capabilities
add_card(slide, Inches(5.8), Inches(1.9), Inches(6.8), Inches(5.0), BG_CARD)
add_accent_bar(slide, Inches(5.8), Inches(1.9), Inches(6.8), Inches(0.06), DT_ORANGE)

pdf_features = [
    ('Company-Specific Branding', 'Auto-populates customer name, industry context,\nand domain-specific terminology', DT_GREEN),
    ('Industry-Contextual ROI', 'Randomised ROI narratives matched to vertical:\nbanking fraud rates, healthcare outcomes, retail conversion', DT_ORANGE),
    ('Journey Step Analysis', 'Maps each customer journey step with\nobservability signals and business metrics', DT_CYAN),
    ('Colour-Coded Visualisations', 'Dynatrace-branded palette with\ncategory-specific colours for each metric type', DT_PURPLE),
    ('Instant Generation', 'POST /api/pdf/executive-summary\nGenerated in seconds — email-ready PDF', DT_BLUE),
]

for j, (title, desc, color) in enumerate(pdf_features):
    y_pos = Inches(2.2 + j * 0.9)
    add_icon_circle(slide, Inches(6.1), y_pos + Inches(0.05), Inches(0.35), color, '', 10)
    add_text_box(slide, Inches(6.6), y_pos, Inches(5.5), Inches(0.3),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(6.6), y_pos + Inches(0.35), Inches(5.8), Inches(0.5),
                 desc, font_size=11, color=GRAY_LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13: DEMO SCENARIOS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 12, 'Four Demo Scenarios You Can Run Tomorrow')

scenarios = [
    {
        'title': 'Zero to Executive Brief',
        'time': '90 SECONDS',
        'steps': [
            'Enter company name + journey type',
            'System auto-generates 6-step journey',
            'BizEvent schema + dashboard created',
            'PDF executive summary generated',
            'Email-ready for C-suite',
        ],
        'color': DT_GREEN,
    },
    {
        'title': 'Industry Vertical Deep Dive',
        'time': '5 MINUTES',
        'steps': [
            'Click Solutions tab',
            'Show Healthcare: Epic/Cerner integrations',
            'Highlight AI clinical decision support KPIs',
            'Show dollar ROI projections',
            'Export industry-specific PDF',
        ],
        'color': DT_BLUE,
    },
    {
        'title': 'Live Service Topology',
        'time': '3 MINUTES',
        'steps': [
            'Launch journey simulation',
            'Switch to Dynatrace Services view',
            'Show Smartscape service flow',
            'Click into distributed trace',
            'Show request-level detail',
        ],
        'color': DT_PURPLE,
    },
    {
        'title': 'Chaos + Business Impact',
        'time': '5 MINUTES',
        'steps': [
            'Activate Chaos Control tab',
            'Inject latency into one service',
            'Show Davis problem detection',
            'Show revenue impact on dashboard',
            'Fix-It agent remediates automatically',
        ],
        'color': DT_ORANGE,
    },
]

for i, scen in enumerate(scenarios):
    x = Inches(0.3 + i * 3.2)
    y = Inches(1.35)
    card_w = Inches(3.05)
    card_h = Inches(5.8)
    
    add_card(slide, x, y, card_w, card_h, BG_CARD)
    add_accent_bar(slide, x, y, card_w, Inches(0.06), scen['color'])
    
    # Scenario number
    add_icon_circle(slide, x + Inches(0.98), y + Inches(0.3), Inches(0.7), scen['color'], str(i+1), 26)
    
    # Title
    add_text_box(slide, x + Inches(0.1), y + Inches(1.2), card_w - Inches(0.2), Inches(0.45),
                 scen['title'], font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Time badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        x + Inches(0.75), y + Inches(1.75), Inches(1.5), Inches(0.3))
    badge.fill.solid()
    badge.fill.fore_color.rgb = scen['color']
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = scen['time']
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    
    # Steps
    for j, step in enumerate(scen['steps']):
        y_step = y + Inches(2.3 + j * 0.55)
        # Step number
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            x + Inches(0.2), y_step + Inches(0.04), Inches(0.22), Inches(0.22))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = BG_CARD_LT
        num_circle.line.color.rgb = scen['color']
        num_circle.line.width = Pt(1)
        tf = num_circle.text_frame
        p = tf.paragraphs[0]
        p.text = str(j + 1)
        p.font.size = Pt(8)
        p.font.color.rgb = scen['color']
        p.alignment = PP_ALIGN.CENTER
        
        add_text_box(slide, x + Inches(0.5), y_step, card_w - Inches(0.7), Inches(0.3),
                     step, font_size=11, color=GRAY_LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14: FAQ / OBJECTION HANDLING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 13, 'Common Questions — Your Answers Ready')

faqs = [
    ('Does this work on Managed?',
     'Absolutely. Core engine works identically — real services, Smartscape, traces, Davis, chaos engineering. '
     'BizEvent querying and AI dashboards are SaaS features. Use the Forge as the SaaS migration conversation.',
     DT_GREEN),
    ('How long does setup take?',
     '15 minutes. Linux VM + Node.js + Dynatrace tenant + API token. Run setup script, answer 6 prompts, you\'re live. '
     'Partners have set this up during a lunch break before an afternoon customer meeting.',
     DT_BLUE),
    ('Can I use this for a POC?',
     'That\'s exactly what it\'s designed for. Pick the customer\'s industry, generate their journey, let it run for a day. '
     'You\'ll have real Dynatrace data that demonstrates production-like observability. The PDF becomes your POC summary.',
     DT_PURPLE),
    ('What about data privacy?',
     'All data is 100% synthetic. Company names, customer profiles, transaction values — everything is generated. '
     'No real customer data involved. Safe for any regulated environment.',
     DT_CYAN),
    ('What if their industry isn\'t covered?',
     'It almost certainly is — 55 verticals from retail banking to space & satellite. But for niche cases, '
     'describe the process in plain English and the AI generates services, schema, and dashboard automatically.',
     DT_ORANGE),
]

for i, (question, answer, color) in enumerate(faqs):
    y_pos = Inches(1.3 + i * 1.15)
    add_card(slide, Inches(0.5), y_pos, Inches(12.3), Inches(1.05), BG_CARD)
    add_accent_bar(slide, Inches(0.5), y_pos, Inches(0.06), Inches(1.05), color)
    
    # Q icon
    q_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), y_pos + Inches(0.15),
                                       Inches(0.35), Inches(0.35))
    q_circle.fill.solid()
    q_circle.fill.fore_color.rgb = color
    q_circle.line.fill.background()
    tf = q_circle.text_frame
    p = tf.paragraphs[0]
    p.text = 'Q'
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, Inches(1.25), y_pos + Inches(0.08), Inches(11), Inches(0.35),
                 f'"{question}"', font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.25), y_pos + Inches(0.48), Inches(11.2), Inches(0.5),
                 answer, font_size=11, color=GRAY_LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15: 90-SECOND POWER DEMO SCRIPT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_header(slide, 14, '90-Second Power Demo — Booth Script', DT_ORANGE)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
             'For hallway conversations, booth demos, and elevator pitches. Memorise this flow.',
             font_size=15, color=GRAY_LIGHT, font_name='Segoe UI Light')

script_steps = [
    {
        'time': '0:00',
        'say': '"Give me an industry — any industry."',
        'do': 'Wait for answer',
        'color': DT_GREEN,
    },
    {
        'time': '0:10',
        'say': '"Watch this — 6 real microservices just spun up."',
        'do': 'Select industry → Launch journey',
        'color': DT_BLUE,
    },
    {
        'time': '0:25',
        'say': '"OneAgent detected them all. Full topology, zero config."',
        'do': 'Show Dynatrace Services + Smartscape',
        'color': DT_CYAN,
    },
    {
        'time': '0:40',
        'say': '"Every transaction carries revenue context — £45K loan value,\n128K lifetime value, 78% conversion rate."',
        'do': 'Show BizEvent data / dashboard',
        'color': DT_PURPLE,
    },
    {
        'time': '1:00',
        'say': '"Now I\'ll break the fraud detection service..."',
        'do': 'Inject chaos → Show Davis problem',
        'color': DT_RED,
    },
    {
        'time': '1:20',
        'say': '"£240K in claims stuck. That\'s the number your CFO cares about.\nThis works on Managed today. 55 industries. Ready to go."',
        'do': 'Show business impact → Close',
        'color': DT_ORANGE,
    },
]

for i, step in enumerate(script_steps):
    y_pos = Inches(1.8 + i * 0.9)
    
    # Time badge
    time_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y_pos, Inches(0.8), Inches(0.35))
    time_badge.fill.solid()
    time_badge.fill.fore_color.rgb = step['color']
    time_badge.line.fill.background()
    tf = time_badge.text_frame
    p = tf.paragraphs[0]
    p.text = step['time']
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    
    # Say column
    add_text_box(slide, Inches(1.5), y_pos - Inches(0.02), Inches(0.5), Inches(0.3),
                 'SAY:', font_size=9, color=step['color'], bold=True)
    add_text_box(slide, Inches(2.0), y_pos - Inches(0.02), Inches(6.5), Inches(0.7),
                 step['say'], font_size=12, color=WHITE)
    
    # Do column
    add_text_box(slide, Inches(8.8), y_pos - Inches(0.02), Inches(0.5), Inches(0.3),
                 'DO:', font_size=9, color=GRAY_MID, bold=True)
    add_text_box(slide, Inches(9.3), y_pos - Inches(0.02), Inches(3.5), Inches(0.7),
                 step['do'], font_size=11, color=GRAY_LIGHT)
    
    # Divider
    if i < len(script_steps) - 1:
        add_accent_bar(slide, Inches(0.5), y_pos + Inches(0.75), Inches(12.3), Inches(0.01), GRAY_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16: CLOSING / CTA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide, RGBColor(0x08, 0x0A, 0x18))

# Decorative shapes
for i in range(5):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.5 + i * 2.5), Inches(5.5 + (i % 2) * 0.5),
        Inches(2), Inches(2))
    s.fill.solid()
    s.fill.fore_color.rgb = DT_GREEN
    s.fill.fore_color.brightness = -0.85
    s.line.fill.background()

add_gradient_bar(slide, 0, Inches(0.05), DT_GREEN)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.5),
             'BUSINESS OBSERVABILITY', font_size=20, color=DT_GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0),
             'FORGE', font_size=64, color=WHITE, bold=True, font_name='Segoe UI Black')

add_accent_bar(slide, Inches(0.8), Inches(2.55), Inches(2), Inches(0.06), DT_GREEN)

# Four takeaways
takeaways = [
    ('Pick any industry', '55 verticals, 256+ templates — your customer\'s industry is covered.', DT_GREEN),
    ('Demo in minutes', 'Real services, real traces, real Smartscape — not PowerPoint.', DT_BLUE),
    ('Leave behind a PDF', 'Executive summary with ROI language that speaks to the C-suite.', DT_PURPLE),
    ('Managed-ready today', 'Core observability works now. SaaS features become the upgrade story.', DT_ORANGE),
]

for i, (title, desc, color) in enumerate(takeaways):
    x = Inches(0.5 + i * 3.15)
    y = Inches(3.2)
    
    add_card(slide, x, y, Inches(3.0), Inches(1.8), BG_CARD)
    add_accent_bar(slide, x, y, Inches(3.0), Inches(0.06), color)
    
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.25),
                                         Inches(0.35), Inches(0.35))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = color
    num_circle.line.fill.background()
    tf = num_circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, x + Inches(0.65), y + Inches(0.22), Inches(2.2), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.6), Inches(0.8),
                 desc, font_size=11, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# Setup CTA
add_card(slide, Inches(2.5), Inches(5.5), Inches(8.3), Inches(1.2), DT_GREEN)
add_text_box(slide, Inches(2.8), Inches(5.6), Inches(7.7), Inches(0.5),
             'GET STARTED', font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.8), Inches(6.1), Inches(7.7), Inches(0.5),
             'Setup takes 15 minutes  |  Linux VM + Node.js + Dynatrace tenant  |  Let\'s talk after this session',
             font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER, font_name='Segoe UI Light')


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           'Business-Observability-Forge-Partner-Event.pptx')
prs.save(output_path)
print(f'\n  PowerPoint saved to: {output_path}')
print(f'  Slides: {len(prs.slides)}')
print(f'  Format: 16:9 widescreen (13.333" x 7.5")')
print(f'\n  Slide Index:')
titles = [
    '1. Title — Business Observability Forge',
    '2. The Partner Challenge (3 problem cards)',
    '3. What Is the Forge? (Architecture)',
    '4. How It Works: Journey Flow + Live Demo',
    '5. Business Events: The Revenue Bridge',
    '6. 55 Industry Verticals (Full Grid)',
    '7. Top 10 Verticals for Partner Conversations',
    '8. Chaos Engineering & AI Remediation',
    '9. Solutions: Healthcare, Banking, Pharma',
    '10. Managed Compatibility Matrix',
    '11. Four AI Agents',
    '12. Executive PDF Export',
    '13. Four Demo Scenarios',
    '14. FAQ / Objection Handling',
    '15. 90-Second Power Demo Script',
    '16. Close & Call to Action',
]
for t in titles:
    print(f'     {t}')
print()
